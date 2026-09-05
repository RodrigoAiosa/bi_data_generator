# -*- coding: utf-8 -*-
"""
scripts/sync_apresentacao.py

Recalcula os números reais do projeto (quantidade de setores, de medidas DAX
geradas, e de abas/ferramentas) e atualiza esses números diretamente dentro
de assets/BI_Data_Generator_Apresentacao.pptx, sem tocar em mais nada do
arquivo (preserva toda a formatação existente).

Rodado automaticamente pelo GitHub Action .github/workflows/sync-apresentacao.yml
a cada push na branch main que mexa em config.py, app.py ou generators/**.
Também pode ser rodado manualmente:

    python scripts/sync_apresentacao.py           # aplica as mudanças
    python scripts/sync_apresentacao.py --check   # só mostra o que mudaria, não salva

IMPORTANTE — o que este script NÃO faz: se uma aba/ferramenta NOVA for
adicionada (ex.: uma 11ª aba), este script atualiza o NÚMERO "10 -> 11" em
todo lugar que ele aparece, mas NÃO cria um card novo no grid de ferramentas
nem um slide dedicado pra essa ferramenta nova — isso exige decisão de
design (texto, layout, o que cabe sem estourar a caixa) e continua sendo
trabalho manual. O script avisa claramente quando detecta esse descompasso
(quantidade de abas != quantidade de cards no grid), pra não deixar isso
passar batido.
"""

from __future__ import annotations

import argparse
import datetime
import re
import sys
import warnings
from pathlib import Path

warnings.filterwarnings("ignore")

RAIZ = Path(__file__).resolve().parent.parent
CAMINHO_PPTX = RAIZ / "assets" / "BI_Data_Generator_Apresentacao.pptx"

sys.path.insert(0, str(RAIZ))


def contar_setores() -> int:
    from config import SETORES
    return len(SETORES)


def contar_medidas_dax(amostra_por_setor: int = 50) -> int:
    """Gera uma amostra pequena de cada setor (rápido, ~5s pros 200) e soma
    a bateria de medidas DAX sugeridas — a mesma lógica usada em produção."""
    from config import SETORES, obter_gerador
    from generators.medidas import gerar_bateria_medidas

    inicio = datetime.date(2024, 1, 1)
    fim = datetime.date(2024, 12, 31)
    total = 0
    for nome in SETORES:
        fn = obter_gerador(nome)
        tabelas = fn(amostra_por_setor, inicio, fim)
        medidas = gerar_bateria_medidas(tabelas)
        total += sum(len(lista) for cats in medidas.values() for lista in cats.values())
    return total


def contar_ferramentas() -> int:
    """Conta quantas abas existem de verdade em app.py, lendo a lista
    passada pra st.tabs([...]) — não um número fixo mantido à mão."""
    texto_app = (RAIZ / "app.py").read_text(encoding="utf-8")
    m = re.search(r"st\.tabs\(\s*\[(.*?)\]\s*,?\s*\)", texto_app, re.DOTALL)
    if not m:
        raise RuntimeError("Não encontrei a chamada st.tabs([...]) em app.py — layout mudou?")
    lista_literal = m.group(1)
    # Cada aba é uma string entre aspas — conta quantas strings tem na lista
    abas = re.findall(r'"[^"]*"', lista_literal)
    return len(abas)


def _set_texto_paragrafo(paragraph, novo_texto: str) -> None:
    runs = paragraph.runs
    if not runs:
        return
    runs[0].text = novo_texto
    for r in runs[1:]:
        r.text = ""


def _texto_completo(shape) -> str:
    if not shape.has_text_frame:
        return ""
    return "".join(r.text for p in shape.text_frame.paragraphs for r in p.runs)


def atualizar_pptx(n_setores: int, n_medidas: int, n_ferramentas: int, aplicar: bool) -> list[str]:
    """Atualiza os 3 números em todo lugar conhecido do pptx. Devolve a
    lista de mudanças feitas (ou que seriam feitas, se aplicar=False)."""
    from pptx import Presentation

    prs = Presentation(str(CAMINHO_PPTX))
    mudancas: list[str] = []

    medidas_fmt = f"{n_medidas:,}".replace(",", ".")

    # ── Slide 1 (índice 0): linha "200 SETORES · 10.696 MEDIDAS DAX · 10 FERRAMENTAS"
    s1 = prs.slides[0]
    for shape in s1.shapes:
        t = _texto_completo(shape)
        if "SETORES" in t.upper() and "FERRAMENTAS" in t.upper():
            nova = f"{n_setores} SETORES  ·  {medidas_fmt} MEDIDAS DAX  ·  {n_ferramentas} FERRAMENTAS"
            if t.strip() != nova.strip():
                mudancas.append(f"Slide 1: {t!r} -> {nova!r}")
                if aplicar:
                    for p in shape.text_frame.paragraphs:
                        if "SETORES" in "".join(r.text for r in p.runs).upper():
                            _set_texto_paragrafo(p, nova)

    # ── Slide 5 (índice 4): "Em números" — pares (número, legenda) por posição X
    s5 = prs.slides[4]
    pares: dict[int, list] = {}
    for shape in s5.shapes:
        if not shape.has_text_frame:
            continue
        pares.setdefault(shape.left, []).append(shape)
    for left, shapes_na_coluna in pares.items():
        if len(shapes_na_coluna) != 2:
            continue
        shapes_na_coluna.sort(key=lambda s: s.top)
        shape_num, shape_legenda = shapes_na_coluna
        legenda = _texto_completo(shape_legenda).lower()
        valor_atual = _texto_completo(shape_num)
        if "setor" in legenda:
            novo_valor = str(n_setores)
        elif "medidas" in legenda:
            novo_valor = medidas_fmt
        elif "ferramentas" in legenda:
            novo_valor = str(n_ferramentas)
        else:
            continue  # "100 mil" (linhas máx.) — não é derivado do código, não mexe
        if valor_atual.strip() != novo_valor.strip():
            mudancas.append(f"Slide 5 ({legenda.strip()}): {valor_atual!r} -> {novo_valor!r}")
            if aplicar:
                _set_texto_paragrafo(shape_num.text_frame.paragraphs[0], novo_valor)

    # ── Slide 7 (índice 6): título "200 setores, modelo estrela pronto"
    s7 = prs.slides[6]
    for shape in s7.shapes:
        t = _texto_completo(shape)
        m = re.match(r"^\d+ setores, (.*)$", t)
        if m:
            nova = f"{n_setores} setores, {m.group(1)}"
            if t.strip() != nova.strip():
                mudancas.append(f"Slide 7: {t!r} -> {nova!r}")
                if aplicar:
                    _set_texto_paragrafo(shape.text_frame.paragraphs[0], nova)

    if aplicar and mudancas:
        prs.save(str(CAMINHO_PPTX))

    return mudancas


def verificar_consistencia_estrutural(n_ferramentas: int) -> str | None:
    """Conta quantos cards numerados (01, 02, ...) existem no grid de
    ferramentas (slide 6) e compara com a quantidade real de abas. Isso NÃO
    é corrigido automaticamente — exige desenho manual de um card/slide novo.
    Devolve uma mensagem de aviso, ou None se estiver tudo batendo."""
    from pptx import Presentation

    prs = Presentation(str(CAMINHO_PPTX))
    s6 = prs.slides[5]
    numeros_no_grid = set()
    for shape in s6.shapes:
        t = _texto_completo(shape).strip()
        if re.fullmatch(r"\d{2}", t):
            numeros_no_grid.add(t)

    if len(numeros_no_grid) != n_ferramentas:
        return (
            f"⚠️  ATENÇÃO: app.py tem {n_ferramentas} aba(s), mas o grid de "
            f"ferramentas no slide 6 tem {len(numeros_no_grid)} card(s) numerado(s). "
            f"Isso o script NÃO corrige sozinho — precisa desenhar manualmente o "
            f"card novo (e possivelmente um slide dedicado pra ferramenta nova), "
            f"do jeito que foi feito pra 'Pergunte aos Dados' e 'Carrossel Power BI'."
        )
    return None


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--check", action="store_true", help="Só mostra o que mudaria, não salva o arquivo.")
    args = parser.parse_args()

    print("Recalculando números reais do projeto...")
    n_setores = contar_setores()
    n_ferramentas = contar_ferramentas()
    print(f"  Setores: {n_setores}")
    print(f"  Ferramentas (abas): {n_ferramentas}")
    print("  Medidas DAX (gerando amostra dos 200 setores, ~5s)...")
    n_medidas = contar_medidas_dax()
    print(f"  Medidas DAX: {n_medidas}")

    mudancas = atualizar_pptx(n_setores, n_medidas, n_ferramentas, aplicar=not args.check)

    if not mudancas:
        print("\n✅ A apresentação já está com os números corretos — nada pra atualizar.")
    else:
        verbo = "Seriam feitas" if args.check else "Foram feitas"
        print(f"\n{verbo} {len(mudancas)} mudança(s):")
        for m in mudancas:
            print(f"  - {m}")
        if not args.check:
            print(f"\n✅ Arquivo salvo: {CAMINHO_PPTX}")

    aviso = verificar_consistencia_estrutural(n_ferramentas)
    if aviso:
        print(f"\n{aviso}")

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
